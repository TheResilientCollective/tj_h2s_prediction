"""Build 3-slide H2S model status deck."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BG        = RGBColor(0x0F, 0x1B, 0x2D)
ACCENT    = RGBColor(0x2E, 0x86, 0xC1)
GREEN     = RGBColor(0x1A, 0xBC, 0x9C)
ORANGE    = RGBColor(0xE6, 0x7E, 0x22)
YELLOW    = RGBColor(0xF1, 0xC4, 0x0F)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xC8, 0xD6, 0xE5)
FIXED_GRN = RGBColor(0x27, 0xAE, 0x60)
DARK_CARD = RGBColor(0x14, 0x26, 0x40)
RULE_DIM  = RGBColor(0x2C, 0x3E, 0x60)
FOOT_COL  = RGBColor(0x6E, 0x8A, 0xA8)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
M   = Inches(0.45)
COL = Inches(4.1)
GAP = Inches(0.22)
CX  = M + COL + GAP
RX  = CX + COL + GAP

# ── primitives ────────────────────────────────────────────────────────────────

def add_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

def rect(slide, l, t, w, h, fill_color, line_color=None, line_width=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
        if line_width:
            s.line.width = Pt(line_width)
    else:
        s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h, size=16, bold=False, italic=False,
       color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box

def bullets(slide, items, l, t, w, h, size=14, color=WHITE, gap=8):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(gap)
        r = p.add_run()
        r.text = f"›  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color

def col_header(slide, text, l, t, w=None, color=ACCENT):
    w = w or COL
    rect(slide, l, t, w, Inches(0.03), color)
    tb(slide, text, l, t + Inches(0.06), w, Inches(0.32),
       size=12, bold=True, color=color)

def slide_header(slide, num, title, subtitle, num_color=ACCENT):
    rect(slide, SLIDE_W - Inches(0.72), Inches(0.16),
         Inches(0.36), Inches(0.36), num_color)
    tb(slide, str(num),
       SLIDE_W - Inches(0.72), Inches(0.16), Inches(0.36), Inches(0.36),
       size=12, bold=True, color=BG if num_color == YELLOW else WHITE,
       align=PP_ALIGN.CENTER)
    rect(slide, M, Inches(0.20), Inches(12.4), Inches(0.04), num_color)
    tb(slide, title, M, Inches(0.28), Inches(9), Inches(0.58),
       size=28, bold=True, color=WHITE)
    tb(slide, subtitle, M, Inches(0.90), Inches(11.5), Inches(0.36),
       size=14, italic=True, color=LIGHT)
    rect(slide, M, Inches(1.32), Inches(12.4), Inches(0.02), RULE_DIM)

def footer(slide):
    rect(slide, M, Inches(7.15), Inches(12.4), Inches(0.02), RULE_DIM)
    tb(slide, "Tijuana River Valley  H₂S Prediction  ·  Model Status  ·  June 2026",
       M, Inches(7.2), Inches(12), Inches(0.28),
       size=10, italic=True, color=FOOT_COL)

# ── presentation ──────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1  —  What we have working
# ═════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
add_bg(s1)
slide_header(s1, 1,
    "What We Have Working",
    "Hourly 3-class XGBoost at NESTOR-BES  ·  Primary production system  ·  Active")
footer(s1)

BODY_TOP = Inches(1.50)

# left: model
col_header(s1, "Model", M, BODY_TOP)
bullets(s1, [
    "Single XGBoost trained on NESTOR-BES observations",
    "43-feature input: weather, wind, flow, H₂S lags, SBIWTP, atmospheric stability",
    "Three output classes: Green (<5 ppb) · Yellow (5–30 ppb) · Orange (≥30 ppb)",
    "hour_sin/cos + stable_atm + is_night — model sees time-of-day signal",
    "Ongoing retraining via Dagster hourly pipeline",
], M, BODY_TOP + Inches(0.42), COL, Inches(3.6), size=15, gap=14)

# centre: performance — one row per metric, no embedded newlines
col_header(s1, "Performance", CX, BODY_TOP)

perf_rows = [
    ("Orange recall",     "61.3 %",                                       GREEN),
    ("False alarm rate",  "5.4 %",                                        FIXED_GRN),
    ("Night (20h–06h)",   "Higher recall — most events are nocturnal",    GREEN),
    ("Day  (06h–20h)",    "Lower recall, consistent FAR — not day-blind", YELLOW),
    ("3-class benefit",   "Yellow captures moderate events a binary split misses", WHITE),
]

box = s1.shapes.add_textbox(CX, BODY_TOP + Inches(0.42), COL, Inches(3.6))
box.word_wrap = True
tf = box.text_frame
tf.word_wrap = True
for i, (key, val, vcol) in enumerate(perf_rows):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_before = Pt(0 if i == 0 else 12)
    rk = p.add_run()
    rk.text = key + ":"
    rk.font.size = Pt(13)
    rk.font.color.rgb = LIGHT

    # value on same paragraph, bolded and coloured
    rv = p.add_run()
    rv.text = "  " + val
    rv.font.size = Pt(14)
    rv.font.bold = True
    rv.font.color.rgb = vcol

# right: limits
col_header(s1, "Known Limits", RX, BODY_TOP, color=ORANGE)
bullets(s1, [
    "NESTOR-BES only — no outputs for IB CIVIC CTR or SAN YSIDRO",
    "No lead-time resolution — same model weights used at 0 h and 24 h ahead",
    "Multi-horizon pipeline is built (36 models, 3 stations, 4 horizons) but STOPPED — needs out-of-season validation before restart",
], RX, BODY_TOP + Inches(0.42), COL, Inches(3.6), size=15, color=LIGHT, gap=18)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2  —  What we fixed
# ═════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
add_bg(s2)
slide_header(s2, 2,
    "What We Fixed",
    "Tiered alert system (T1–T3)  ·  Complete collapse Apr–May 2026  →  Repaired 2026-06-02",
    num_color=FIXED_GRN)
footer(s2)

# ── left: collapse table + root causes ───────────────────────────────────────
col_header(s2, "The Collapse", M, BODY_TOP, color=ORANGE)

collapse_data = [
    ("Month",      "% fired", "Real H₂S?"),
    ("Feb 2026",   "5.5 %",   "yes"),
    ("Mar 2026",   "15.5 %",  "yes"),
    ("Apr 2026 ▼", "0.0 %",   "48 ppb avg"),
    ("May 2026 ▼", "0.0 %",   "13 ppb avg"),
]
CW = [Inches(1.35), Inches(1.1), Inches(1.25)]
row_h = Inches(0.37)
ty = BODY_TOP + Inches(0.42)
for ri, row in enumerate(collapse_data):
    cx2 = M
    for ci, cell in enumerate(row):
        is_hdr  = ri == 0
        is_dead = ri >= 3 and ci in (1, 2)
        rect(s2, cx2, ty, CW[ci], row_h, ACCENT if is_hdr else DARK_CARD)
        tb(s2, cell, cx2 + Inches(0.05), ty + Inches(0.05),
           CW[ci] - Inches(0.07), row_h - Inches(0.06),
           size=12, bold=is_hdr or is_dead,
           color=ORANGE if is_dead else WHITE,
           align=PP_ALIGN.CENTER)
        cx2 += CW[ci]
    ty += row_h

tb(s2, "Root causes", M, ty + Inches(0.18), COL, Inches(0.28),
   size=12, bold=True, color=LIGHT)
bullets(s2, [
    "SBIWTP hard gate: plant at surplus all spring → anomaly ≥ 0 → entire alert chain dead",
    "SBIWTP weights 5–9× too large → surplus pushed score down",
    "quiet_night_stats were placeholders — sbiwtp_anomaly std 22× inflated",
], M, ty + Inches(0.50), COL, Inches(2.2), size=13, color=LIGHT, gap=8)

# ── centre: the four fixes ────────────────────────────────────────────────────
col_header(s2, "The Fixes  (v4 config, 2026-06-02)", CX, BODY_TOP, color=FIXED_GRN)
bullets(s2, [
    "Gate: SBIWTP deficit prerequisite replaced by met-data availability (wind_speed_10m not NaN)",
    "Weights re-derived from Oct 2024–Jun 2026 full dataset (Cohen's d):\n"
     "  sbiwtp_flow  −1.44 → −0.24\n"
     "  flow_log sign flip → +0.38\n"
     "  stable_atm → +0.50   |   wind_speed remains −0.54",
    "Tier 1 expanded: atmospheric features added so calm/stable nights unlock the chain even with no SBIWTP deficit",
    "quiet_night_stats calibrated from 1,230 real quiet-night rows — sbiwtp_anomaly std 3.5 → 0.16",
], CX, BODY_TOP + Inches(0.42), COL, Inches(4.8), size=13, gap=12)

# ── right: before/after recall table ─────────────────────────────────────────
col_header(s2, "Recall  Before → After  (Tier 3)", RX, BODY_TOP, color=ACCENT)

ba_data = [
    ("Month",      "Horizon",    "Before", "After"),
    ("Apr 2026",   "nowcast",    "0.08",   "0.54"),
    ("Apr 2026",   "mid",        "0.08",   "0.45"),
    ("Apr 2026",   "day-ahead",  "0.07",   "0.29"),
    ("May 2026",   "nowcast",    "0.00",   "0.48"),
    ("May 2026",   "mid",        "0.00",   "0.41"),
    ("May 2026",   "day-ahead",  "0.00",   "0.24"),
    ("All months", "nowcast",    "0.23",   "0.35"),
]
BCW = [Inches(1.15), Inches(1.05), Inches(0.77), Inches(0.83)]
row_h2 = Inches(0.40)
ty2 = BODY_TOP + Inches(0.42)
for ri, row in enumerate(ba_data):
    cx3 = RX
    for ci, cell in enumerate(row):
        is_hdr  = ri == 0
        is_aftr = ci == 3 and ri > 0
        is_bfr  = ci == 2 and ri > 0
        rect(s2, cx3, ty2, BCW[ci], row_h2, ACCENT if is_hdr else DARK_CARD)
        tb(s2, cell, cx3 + Inches(0.04), ty2 + Inches(0.05),
           BCW[ci] - Inches(0.05), row_h2 - Inches(0.06),
           size=12, bold=is_hdr or is_aftr,
           color=(FIXED_GRN if is_aftr else ORANGE if is_bfr else WHITE),
           align=PP_ALIGN.CENTER)
        cx3 += BCW[ci]
    ty2 += row_h2

tb(s2, "System no longer collapses on atmospheric-driven events.",
   RX, ty2 + Inches(0.12), COL, Inches(0.4),
   size=12, bold=True, color=FIXED_GRN)

# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3  —  What's left
# ═════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
add_bg(s3)
slide_header(s3, 3,
    "What's Left",
    "Remaining gaps after the fix  ·  Open work items  ·  Path to broader coverage",
    num_color=YELLOW)
footer(s3)

CARD_W = Inches(3.9)
CARD_H = Inches(5.1)
CARD_T = Inches(1.42)

card_defs = [
    ("Gap 1  —  Daytime alerting", ORANGE, WHITE, M, [
        "Zero fires on daytime hours in Apr/May 2026 — fix did not help daytime",
        "Sea breeze raises wind speed above 4.0 m/s gate during daylight — gate never passes",
        "Daytime H₂S events driven by different mechanisms than nocturnal",
        "Needs separate daytime gate thresholds or a parallel daytime scoring path",
        "Left as future work — requires labeled daytime event dataset to calibrate",
    ]),
    ("Gap 2  —  Tiered system maturity", YELLOW, WHITE, CX, [
        "Design acceptance targets (§6.1) not yet met",
        "Need ≥2 years of labeled events for robust threshold tuning",
        "quiet_night_stats must be regenerated after each dataset refresh\n(uv run python -m h2s.defs.tiered_alerts.backtest --emit-stats)\n→ add to operational runbook",
        "NB fallback to IB CIVIC CTR implemented but unvalidated during outages",
    ]),
    ("Gap 3  —  Single-station scope", ACCENT, WHITE, RX, [
        "Hourly model: NESTOR-BES only — no IB CIVIC CTR or SAN YSIDRO",
        "No lead-time resolution — 0 h and 24 h use identical weights",
        "Multi-horizon pipeline built: 36 models, 3 stations, 4 horizon buckets",
        "STOPPED — trained on Nov 2023–Apr 2026 only, no out-of-season validation",
        "Path: accumulate May–Oct 2026 data, validate seasonal generalization, then restart",
    ]),
]

for label, border_color, label_text_color, cx_pos, body in card_defs:
    rect(s3, cx_pos, CARD_T, CARD_W, CARD_H, DARK_CARD,
         line_color=border_color, line_width=1.5)
    rect(s3, cx_pos, CARD_T, CARD_W, Inches(0.06), border_color)
    tb(s3, label,
       cx_pos + Inches(0.12), CARD_T + Inches(0.10),
       CARD_W - Inches(0.22), Inches(0.38),
       size=12, bold=True, color=label_text_color)
    rect(s3, cx_pos + Inches(0.12), CARD_T + Inches(0.54),
         CARD_W - Inches(0.24), Inches(0.02), RULE_DIM)
    bullets(s3, body,
            cx_pos + Inches(0.12), CARD_T + Inches(0.62),
            CARD_W - Inches(0.22), CARD_H - Inches(0.78),
            size=13, color=LIGHT, gap=7)

# ── save ──────────────────────────────────────────────────────────────────────
out = "/Users/valentin/development/dev_resilient/tj_h2s_prediction/docs/h2s_new_model/h2s_model_status_june2026.pptx"
prs.save(out)
print(f"Saved: {out}")
